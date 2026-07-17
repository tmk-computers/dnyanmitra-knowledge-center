import sys
import os
import urllib.parse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_footer(slide):
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4))
    tf = footer_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Confidential | DASP Digital Pvt Ltd | www.dnyanmitra.com"
    p.font.name = 'Calibri'
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title, subtitle, presenter):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 60, 120) # Blue
    bg.line.fill.background()
    
    bar = slide.shapes.add_shape(1, Inches(0.0), Inches(7.2), Inches(13.333), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(245, 130, 32) # Orange
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Calibri'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.333), Inches(1.2))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = 'Calibri'
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(245, 130, 32)
    
    pres_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.2))
    tf3 = pres_box.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = f"Presenter: {presenter}\nDASP Digital Pvt Ltd\nTerritory: Haveli Taluka, Pune"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(200, 200, 200)
    
    slide.notes_slide.notes_text_frame.text = "Welcome and cover slide. Introduce the session as the DnyanMitra Taluka Head Due Diligence briefing for Haveli Taluka, Pune."

def add_section_divider(prs, section_num, section_title):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 60, 120)
    bg.line.fill.background()
    
    sidebar = slide.shapes.add_shape(1, Inches(0.0), Inches(0.0), Inches(0.4), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = RGBColor(245, 130, 32)
    sidebar.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.5), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"SECTION {section_num}"
    p.font.name = 'Calibri'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(245, 130, 32)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = section_title
    p2.font.name = 'Calibri'
    p2.font.size = Pt(36)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.font.bold = True
    p2.space_before = Pt(10)
    
    slide.notes_slide.notes_text_frame.text = f"Transitioning to Section {section_num}: {section_title}."

def add_content_slide(prs, section_name, title_text, bullets, notes_text):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9))
    tf = header_box.text_frame
    tf.word_wrap = True
    
    p_sec = tf.paragraphs[0]
    p_sec.text = section_name.upper()
    p_sec.font.name = 'Calibri'
    p_sec.font.size = Pt(11)
    p_sec.font.bold = True
    p_sec.font.color.rgb = RGBColor(245, 130, 32)
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.name = 'Calibri'
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(20, 60, 120)
    p_title.space_before = Pt(4)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(245, 130, 32)
    line.line.fill.background()
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.0))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p_bullet = tf_content.paragraphs[0]
        else:
            p_bullet = tf_content.add_paragraph()
        
        is_sub = bullet.startswith("  ")
        cleaned_bullet = bullet.strip()
        
        words = cleaned_bullet.split()
        if len(words) > 8:
            cleaned_bullet = " ".join(words[:8])
            
        p_bullet.text = cleaned_bullet
        p_bullet.font.name = 'Calibri'
        p_bullet.font.size = Pt(15 if is_sub else 17)
        p_bullet.font.color.rgb = RGBColor(70, 70, 70)
        p_bullet.space_after = Pt(8 if is_sub else 12)
        if is_sub:
            p_bullet.level = 1
        else:
            p_bullet.level = 0
            
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = notes_text

def add_table_slide(prs, section_name, title_text, headers, rows, notes_text):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9))
    tf = header_box.text_frame
    tf.word_wrap = True
    
    p_sec = tf.paragraphs[0]
    p_sec.text = section_name.upper()
    p_sec.font.name = 'Calibri'
    p_sec.font.size = Pt(11)
    p_sec.font.bold = True
    p_sec.font.color.rgb = RGBColor(245, 130, 32)
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.name = 'Calibri'
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(20, 60, 120)
    p_title.space_before = Pt(4)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(245, 130, 32)
    line.line.fill.background()
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = Inches(12.333)
    height = Inches(4.8)
    
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(20, 60, 120)
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(70, 70, 70)
            p.alignment = PP_ALIGN.LEFT
            
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = notes_text

def add_roadmap_slide(prs, section_name, title_text, milestones, notes_text):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9))
    tf = header_box.text_frame
    tf.word_wrap = True
    
    p_sec = tf.paragraphs[0]
    p_sec.text = section_name.upper()
    p_sec.font.name = 'Calibri'
    p_sec.font.size = Pt(11)
    p_sec.font.bold = True
    p_sec.font.color.rgb = RGBColor(245, 130, 32)
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.name = 'Calibri'
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(20, 60, 120)
    p_title.space_before = Pt(4)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(245, 130, 32)
    line.line.fill.background()
    
    timeline_y = Inches(3.6)
    t_line = slide.shapes.add_shape(1, Inches(1.0), timeline_y, Inches(11.333), Inches(0.06))
    t_line.fill.solid()
    t_line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    t_line.line.fill.background()
    
    num_steps = len(milestones)
    step_width = Inches(11.333) / num_steps
    
    for i, ms in enumerate(milestones):
        x = Inches(1.0) + (i * step_width) + (step_width / 2) - Inches(1.0)
        
        circle = slide.shapes.add_shape(9, x + Inches(0.85), timeline_y - Inches(0.12), Inches(0.24), Inches(0.24))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(245, 130, 32)
        circle.line.fill.background()
        
        box = slide.shapes.add_shape(1, x, timeline_y + Inches(0.4), Inches(2.0), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 247, 250)
        box.line.color.rgb = RGBColor(20, 60, 120)
        
        tf_box = box.text_frame
        tf_box.word_wrap = True
        p_t = tf_box.paragraphs[0]
        p_t.text = ms['title']
        p_t.font.name = 'Calibri'
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = RGBColor(20, 60, 120)
        p_t.alignment = PP_ALIGN.CENTER
        
        p_d = tf_box.add_paragraph()
        p_d.text = ms['desc']
        p_d.font.name = 'Calibri'
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = RGBColor(80, 80, 80)
        p_d.alignment = PP_ALIGN.CENTER
        p_d.space_before = Pt(4)
        
        box_deliv = slide.shapes.add_textbox(x, timeline_y - Inches(1.6), Inches(2.0), Inches(1.2))
        tf_deliv = box_deliv.text_frame
        tf_deliv.word_wrap = True
        p_del = tf_deliv.paragraphs[0]
        p_del.text = f"Deliverable:\n{ms['kpi']}"
        p_del.font.name = 'Calibri'
        p_del.font.size = Pt(11)
        p_del.font.bold = True
        p_del.font.color.rgb = RGBColor(46, 125, 50)
        p_del.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = notes_text

def build_presentation(file_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Cover
    add_title_slide(prs, "DnyanMitra\nTaluka Head Due Diligence Kit", "Haveli Taluka, Pune District\nMaharashtra, India", "District Head")
    
    # Slide 2: Agenda
    agenda_bullets = [
        "1. Welcome & Session Objectives",
        "2. About DASP Digital Core",
        "3. DnyanMitra Platform Sourcing Ecosystem",
        "4. Why Sourcing Opportunities Exist",
        "5. Haveli Taluka Regional Sizing Report",
        "6. Taluka Head Role & KPIs",
        "7. Worked Commission Models & Revenue Sharing",
        "8. Detailed Product Catalog"
    ]
    add_content_slide(prs, "Agenda", "Session Brief & Structure Overview", agenda_bullets, "Brief the candidate on the structure of the workshop. Note that sections 9 through 16 cover operational steps, business charters, 12-month calendar plans, risk assessment sheets, and onboarding checklists.")
    
    # SECTION 2
    add_section_divider(prs, 2, "About DASP Digital")
    
    add_content_slide(prs, "About DASP Digital", "Corporate Overview", [
        "Pune-based EdTech and SportsTech company.",
        "Developer of specialized regional marketplace ecosystems.",
        "Operates B2B sourcing frameworks across Maharashtra.",
        "Establishment of Digital Transformation Centres locally.",
        "Zero middle-tier distribution sourcing networks.",
        "Built on absolute financial transaction transparency."
    ], "Introduce DASP Digital as a professional entity based in Pune, outlining its position in the B2B regional digital marketplace sector.")

    add_content_slide(prs, "About DASP Digital", "Corporate Vision & Core Focus", [
        "Pioneering digital transformation for regional institutions.",
        "Targeting tier-2 and tier-3 talukas.",
        "Establishment of localized business transformation networks.",
        "Fostering entrepreneurship via regional partner appointments.",
        "Aggregating unorganized localized hardware merchant networks.",
        "Standardizing delivery quality via physical site verification."
    ], "Define DASP Digital's long-term vision of setting up localized transformation nodes in semi-urban and rural centers.")

    add_content_slide(prs, "About DASP Digital", "Company Mission Statement", [
        "To digitalise local institutional procurement chains.",
        "To reduce sourcing overheads for trustees.",
        "To verify vendor credentials physically and verify banking.",
        "To deploy advanced AI learning modules locally.",
        "To establish trusted technology sourcing channels.",
        "To build sustainable, long-term business partnerships."
    ], "State the core mission of bridging the tech infrastructure gap for schools and PE sports programs outside major urban zones.")

    add_content_slide(prs, "About DASP Digital", "Technology Expertise & Architecture", [
        "AWS hosting with secure regional Mumbai clustering.",
        "Node.js microservices with PostgreSQL transaction layers.",
        "Python-based AI engines for school auditing.",
        "Offline-first mobile-responsive CRM data logging integrations.",
        "Automated WhatsApp notification delivery system interfaces.",
        "Data encryption standard protecting institutional files."
    ], "Discuss DASP Digital's technical infrastructure, focusing on database reliability, AI modules, and secure client cloud databases.")

    add_content_slide(prs, "About DASP Digital", "Platform Suite Overview", [
        "DnyanMitra B2B education procurement marketplace.",
        "KridaMitra sports PE booking network platform.",
        "School and College Cloud ERP softwares.",
        "AI-assisted exam grading automated system portals.",
        "Institutional Excellence (IX) benchmarking scorecards.",
        "Annual Maintenance Contract (AMC) tracker modules."
    ], "Outline the full suite of software and marketplace modules under the DASP Digital umbrella.")

    add_content_slide(prs, "About DASP Digital", "Services Suite Sourcing", [
        "Smart Board hardware installation and alignment.",
        "On-site school computer laboratory engineering setup.",
        "CCTV security mapping and biometric configurations.",
        "On-site Institutional Excellence digital technology audits.",
        "Bi-annual faculty ERP software training workshops.",
        "Preventive quarterly maintenance field engineer dispatches."
    ], "Provide an overview of the technical and consulting services offered through the local Digital Transformation Centres.")

    add_content_slide(prs, "About DASP Digital", "Future Platform Expansion Roadmap", [
        "Year 1: Consolidate Western Maharashtra DTCs.",
        "Year 2: Launch KrushiMitra agricultural marketplace.",
        "Year 3: Deploy localized AI auditing dashboards.",
        "Year 4: Launch ZilaValley pilot procurement.",
        "Year 5: Inter-state expansion to adjacent states.",
        "Platform modularity ensures seamless feature releases."
    ], "Present DASP Digital's 5-year scaling plan, highlighting the addition of agriculture (KrushiMitra) and municipal procurement (ZilaValley).")

    add_content_slide(prs, "About DASP Digital", "Organizational Hierarchy Structure", [
        "Pune HQ Board sets policy standards.",
        "Division Heads oversee regional divisions.",
        "District Heads mentor local Taluka Heads.",
        "Taluka Heads manage individual taluka DTCs.",
        "Field Sales Executives handle daily visits.",
        "Ensures regional support and rapid ticket resolution."
    ], "Brief the candidate on the support network, explaining how they report to District Heads and lead local sales executives.")

    # SECTION 3
    add_section_divider(prs, 3, "About DnyanMitra")
    
    add_content_slide(prs, "About DnyanMitra", "Core Marketplace Concept", [
        "B2B procurement marketplace for educational institutions.",
        "Direct manufacturer links bypass retail middle-tiers.",
        "Vetted vendors ensure delivery safety standard compliance.",
        "Integrated billing system with transparent margin commissions.",
        "Physical product trials hosted at Taluka DTCs.",
        "Covers software licenses, hardware, and physical goods."
    ], "Introduce the DnyanMitra platform as the unified sourcing hub for schools, linking digital ERPs with physical school hardware.")

    add_content_slide(prs, "About DnyanMitra", "Digital Transformation Centre (DTC)", [
        "DTC serves as physical demonstration office.",
        "Showcases interactive smart boards and projectors.",
        "Maintains local stocks of critical replacement parts.",
        "Base office for Field Sales Executives.",
        "Owned and managed by the Taluka Head.",
        "Ensures local presence and builds trustee trust."
    ], "Explain the role of the DTC physical office in the taluka, serving as the local point-of-presence for product testing.")

    add_content_slide(prs, "About DnyanMitra", "Institutional Excellence Framework", [
        "AI-driven framework audits school technology gaps.",
        "Scores schools on safety, admin, pedagogy.",
        "Generates customized digital upgrade audit reports.",
        "Provides school boards with multi-year roadmap.",
        "Drives customer loyalty through structured upgrades.",
        "Audit tool used as primary sales hook."
    ], "Detail the Institutional Excellence (IX) framework and how it acts as a consultative auditing tool for sales pipelines.")

    add_content_slide(prs, "About DnyanMitra", "AI for Education Strategy", [
        "Automated exam paper scans speed up grading.",
        "Predictive vendor sourcing agents match lowest bids.",
        "Smart scheduling algorithms optimize school computer labs.",
        "Student attendance biometrics trigger automated warnings.",
        "Personalized learning suggestions based on student scores.",
        "Provides rural institutions with Tier-1 AI tools."
    ], "Describe DnyanMitra's AI solutions, focusing on smart bidding, automated grading, and diagnostic parent alert gateways.")

    add_content_slide(prs, "About DnyanMitra", "KridaMitra Sports Sourcing", [
        "Integrates physical education into school platforms.",
        "Aggregates sports gear procurement and uniform bids.",
        "Enables automated scheduling for inter-school sports matches.",
        "Scores student physical health metrics via ERP.",
        "Onboards local sports academies and PE coaches.",
        "Unlocks revenue via specialized athletic product sourcing."
    ], "Explain how the KridaMitra module complements school software by addressing athletic equipment sourcing and facility bookings.")

    add_content_slide(prs, "About DnyanMitra", "Stakeholder Ecosystem Matrix", [
        "Trustees get cost savings and compliance.",
        "Principals get administrative dashboards and trackers.",
        "Vendors get direct institutional sales channels.",
        "Taluka Heads earn recurring margin commissions.",
        "FSEs earn direct target-based monthly incentives.",
        "Students get modern smart classroom environments."
    ], "Detail how all participants benefit, establishing DnyanMitra as a collaborative ecosystem, not just a retail site.")

    # SECTION 4
    add_section_divider(prs, 4, "Why This Opportunity Exists")
    
    add_content_slide(prs, "Why DnyanMitra Exists", "School Administrative Problems", [
        "Trustees spend hours managing multiple stationery vendors.",
        "Manual fee registers lead to collection delays.",
        "School buses lack real-time GPS tracking safety.",
        "Fragmented software products do not sync data.",
        "No local hardware support when devices fail.",
        "Staff spend hours processing manual grade books."
    ], "Address the pain points of school principals and trustees, highlighting inefficiencies in manual data and vendor fragmentation.")

    add_content_slide(prs, "Why DnyanMitra Exists", "Sourcing & Procurement Hurdles", [
        "Wholesalers add high markups to school goods.",
        "No standard technical specifications for classroom desks.",
        "Textbook delivery delays disrupt academic year starts.",
        "Unvetted vendors deliver sub-standard school uniforms.",
        "Non-transparent bidding leads to trustee budget waste.",
        "Lack of direct factory lines for rural areas."
    ], "Outline the supply chain inefficiencies in traditional school procurement that DnyanMitra resolves.")

    add_content_slide(prs, "Why DnyanMitra Exists", "Lack of Local Technical Support", [
        "Pune/Mumbai software firms ignore taluka institutions.",
        "Smart boards sit unused due to training gaps.",
        "Hardware repair delays halt classroom smart learning.",
        "No local technicians to verify wiring setups.",
        "High cost of annual maintenance contract options.",
        "DTC model solves support via local presence."
    ], "Explain the support gap that rural school heads face, positioning the local DTC as the primary solution.")

    add_content_slide(prs, "Why DnyanMitra Exists", "The Local Presence Advantage", [
        "Principals prefer dealing with local business leaders.",
        "Physical DTC demos remove technology trust barriers.",
        "Rapid on-site installation support matches target schedules.",
        "Local vendor verification protects marketplace transaction safety.",
        "Regional representation builds long-term board relationships.",
        "DASP provides the platform, partner provides leadership."
    ], "Summarize why the Taluka Head partnership model beats distant, corporate sales operations in semi-urban India.")

    # SECTION 5
    add_section_divider(prs, 5, "Haveli Taluka Opportunity Report")
    
    add_content_slide(prs, "Haveli Taluka Report", "Territory Profile & Geography", [
        "Territory: Haveli Taluka, Pune District, Maharashtra.",
        "Surrounds Pune Municipal Corporation suburban borders.",
        "Contains major clusters: Hadapsar, Loni Kalbhor, Wagholi.",
        "Verified Population: Approximately 24.3 Lakh residents.",
        "Verified Literacy Rate: 88.2% (exceeds state average).",
        "Includes large industrial and educational transition zones."
    ], "Introduce Haveli Taluka as a premium, high-density territory surrounding Pune city, making it a highly profitable DTC market.")

    add_content_slide(prs, "Haveli Taluka Report", "Educational Institution Sizing", [
        "Verified schools (State Board/CBSE): 340+ campuses.",
        "Estimated Senior & Junior Colleges: 45+ institutions.",
        "Estimated ITIs & Polytechnic classes: 12 campuses.",
        "Estimated private coaching academies: 180+ classes.",
        "Estimated student population: Over 1.8 Lakh students.",
        "High concentration of CBSE English-medium boarding schools."
    ], "Provide the institutional counts for Haveli Taluka, highlighting the large target market size for ERP and hardware sales.")

    add_content_slide(prs, "Haveli Taluka Report", "Local Sourcing Vendor Sizing", [
        "Estimated computer & IT dealers: 65 shops.",
        "Estimated uniform and textbook printers: 18 units.",
        "Estimated school furniture fabricators: 12 units.",
        "Estimated electrical & CCTV contractors: 30 units.",
        "Offers massive potential for marketplace vendor onboarding.",
        "DTC can aggregate these into verified local suppliers."
    ], "Provide the local supplier landscape data, demonstrating opportunities for the Taluka Head to onboard local merchants.")

    add_content_slide(prs, "Haveli Taluka Report", "Technology Adoption & Gap Analysis", [
        "High tech adoption in Wagholi/Hadapsar clusters.",
        "Medium adoption in Loni Kalbhor classes.",
        "Low adoption in rural eastern villages.",
        "Verified gap: 60% schools lack smart boards.",
        "Verified gap: 70% primary schools use paper registers.",
        "Creates immediate demand for ERP and classroom devices."
    ], "Discuss the digital readiness of Haveli schools, highlighting the gap in smart learning tools and paperless admin.")

    add_content_slide(prs, "Haveli Taluka Report", "Priority School Clusters", [
        "Cluster 1: Wagholi (High private CBSE density).",
        "Cluster 2: Loni Kalbhor (Colleges, ITIs).",
        "Cluster 3: Hadapsar Suburban (Large private trusts).",
        "DTC office location recommended near Loni Kalbhor.",
        "Ensures rapid road access to eastern rural schools.",
        "Provides central hub for FSE daily school routes."
    ], "Detail the geographic zones to prioritize, suggesting the optimal location for the partner's DTC office.")

    # SWOT Table
    swot_headers = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    swot_rows = [[
        "High school density near Pune city limits.",
        "FSE recruitment costs are higher in Pune.",
        "Smart board demand driven by CBSE rules.",
        "Cheap standalone ERP packages from local shops."
    ], [
        "Well-developed road access across Haveli.",
        "Trustees expect Pune-level support speeds.",
        "Local vendor onboarding builds regional support.",
        "Delays in school trustee board approvals."
    ]]
    add_table_slide(prs, "Haveli Taluka Report", "Haveli Territory SWOT Matrix", swot_headers, swot_rows, "Present a realistic SWOT analysis for Haveli Taluka, highlighting local advantages and regional mitigation steps.")

    # Market size table
    mkt_headers = ["Product Line", "Total Market Sizing", "DTC Year 1 Target Sourcing", "Est. Commission (₹)"]
    mkt_rows = [
        ["School Cloud ERP", "340 Schools (₹68L)", "25 Schools (₹5L)", "₹50,000 (10%)"],
        ["Smart Classrooms", "1,200 Rooms (₹11.4Cr)", "40 Rooms (₹38L)", "₹1,52,000 (4%)"],
        ["Biometrics/CCTV", "300 Setups (₹60L)", "15 Setups (₹3L)", "₹12,000 (4%)"],
        ["Marketplace Sourcing", "₹8.0 Crore Sourcing", "₹20 Lakh Sourcing", "₹40,000 (2%)"],
        ["AMC Agreements", "₹1.2 Crore Contracts", "₹5 Lakh Contracts", "₹60,000 (12%)"]
    ]
    add_table_slide(prs, "Haveli Taluka Report", "Estimated Market Sizing & Projections", mkt_headers, mkt_rows, "Detail the financial projections. Clearly distinguish between total market sizing (estimates) and Year 1 DTC targets.")

    add_content_slide(prs, "Haveli Taluka Report", "DTC Year 1 Sourcing Strategy", [
        "Focus on Wagholi/Hadapsar private schools first.",
        "Onboard 10 premium local IT hardware vendors.",
        "Deploy free initial digital audits as hooks.",
        "Leverage training institute network to recruit FSE.",
        "Organize 1 principal seminar in Loni Kalbhor.",
        "Build recurring software renewal base in Q4."
    ], "Summarize the execution strategy for Haveli Taluka, emphasizing early targets and local resource allocation.")

    # SECTION 6
    add_section_divider(prs, 6, "Taluka Head Role & KPIs")
    
    add_content_slide(prs, "Taluka Head Role", "Core Mission & Leadership Scope", [
        "To establish DASP Digital's DTC in Haveli.",
        "To manage regional school and vendor relations.",
        "To recruit, train, and guide field staff.",
        "To lead the local digital transformation audits.",
        "To represent DASP brands at district meetings.",
        "To build a sustainable regional business partnership."
    ], "Define the leadership role. Emphasize that the Taluka Head is an operational partner and entrepreneur, not a simple sales rep.")

    add_content_slide(prs, "Taluka Head Role", "Key Responsibilities", [
        "Scheduling smart classroom and ERP software demos.",
        "Recruiting and managing one Field Sales Executive.",
        "Conducting physical audits of local vendor premises.",
        "Updating CRM lead tracking logs daily.",
        "Verifying bank accounts and KYC documents physically.",
        "Coordinating with the Pune HQ support desk."
    ], "Detail the core operational tasks required, highlighting team management and physical KYC validation duties.")

    # Cadence table
    cad_headers = ["Cadence", "Task Scope", "Deliverable Checked", "Frequency"]
    cad_rows = [
        ["Daily", "Morning briefing with FSE", "CRM visit routes checked", "Daily at 9:00 AM"],
        ["Daily", "School board/trustee visits", "1-2 physical meetings logged", "Daily PM"],
        ["Weekly", "Vendor onboarding reviews", "2 local vendors audited", "Weekly"],
        ["Weekly", "DH Review Meetings", "Weekly pipeline status report", "Saturdays"],
        ["Monthly", "Target & Commission Audits", "Accounts reconciliation sheets", "Monthly"]
    ]
    add_table_slide(prs, "Taluka Head Role", "DTC Standard Work Cadence", cad_headers, cad_rows, "Walk through the daily, weekly, and monthly tasks to show the structured routine required to run the franchise.")

    add_content_slide(prs, "Taluka Head Role", "Key Performance Indicators (KPIs)", [
        "Active Institution Sourcing Penetration (30% weight).",
        "Field Sales Executive (FSE) Productivity (30% weight).",
        "Onboarded Vendor Quality and Compliance (20% weight).",
        "Platform Transaction Volume processed (20% weight).",
        "Weekly CRM visit logs completeness is mandatory.",
        "Quarterly ratings determine territory contract renewal."
    ], "Explain how the partner's performance is graded. Emphasize that reviews are designed to support growth and resolve hurdles.")

    add_content_slide(prs, "Taluka Head Role", "Support Provided by DASP Digital", [
        "Dedicated District Head mentor for field deals.",
        "Technical deployment engineers for smart classroom installs.",
        "Marketing brochures, banners, and standee print files.",
        "Continuous product feature and ERP module updates.",
        "National vendor procurement contracts and price sheets.",
        "Weekly online billing reconciliation and settlements."
    ], "Reassure the candidate of DASP Digital's backing, detailing the support channels, tech setup, and marketing materials.")

    add_content_slide(prs, "Taluka Head Role", "Escalation Matrix", [
        "L1: Taluka Head DTC (First contact).",
        "L2: District Head (Unresolved support/quoting questions).",
        "L3: Division Head (Contract/territory disputes).",
        "L4: Pune HQ Support Desk (Software/billing bugs).",
        "SLAs require L1 response within 4 hours.",
        "Ensures excellent customer service for local schools."
    ], "Detail the hierarchy for resolving software bugs, billing disputes, and shipping issues.")

    # SECTION 7
    add_section_divider(prs, 7, "Revenue Sharing")
    
    add_content_slide(prs, "Revenue Sharing", "Compensation Philosophy", [
        "Builds mutual commitment through commission-based margins.",
        "Multiple revenue streams ensure steady monthly income.",
        "Zero franchise fees ensures low barrier entry.",
        "Recurring commissions on software renewals reward quality.",
        "Commissions paid directly via monthly bank transfer.",
        "Clear worked examples prevent misunderstandings later."
    ], "Explain the partnership financial mindset, focusing on low barriers, multiple transaction lines, and direct monthly settlements.")

    add_content_slide(prs, "Revenue Sharing", "Sourcing Split Matrix", [
        "Gross Platforms Margin: 100% of distribution pool.",
        "DASP Digital Platform Share: 40% (covers cloud/SaaS).",
        "State/Division Heads Override Share: 15% (oversight).",
        "District Head Override Share: 15% (field training).",
        "Taluka Head DTC Direct Share: 30% (local DTC).",
        "Rewards local execution and physical relationship presence."
    ], "Detail the exact margin split of DASP Digital platform commissions across the regional tiers.")

    # Revenue table
    rev_headers = ["Product Category", "Platform Margin", "Taluka DTC Direct Share", "Worked Example Payout (₹)"]
    rev_rows = [
        ["Software subscriptions", "25% to 30%", "8% to 10%", "₹10,000 on ₹1L ERP Contract"],
        ["Classroom Hardware", "8% to 12%", "3% to 4%", "₹16,000 on ₹4L Smart Board order"],
        ["Marketplace Goods", "5% to 10%", "1.5% to 2.5%", "₹16,000 on ₹8L Uniform supply"],
        ["Maintenance (AMCs)", "25% to 35%", "10% to 12%", "₹6,000 on ₹50K AMC Contract"]
    ]
    add_table_slide(prs, "Revenue Sharing", "Commission Structure by Product Category", rev_headers, rev_rows, "Walk through the commission rates and payouts for software, hardware, bulk goods, and maintenance packages.")

    add_content_slide(prs, "Revenue Sharing", "Worked Example: School ERP Deal", [
        "Client: Wagholi English School (700 students).",
        "ERP SaaS Price: ₹200 per student/year.",
        "Total Annual Invoice Value: ₹1,40,000 (excl. tax).",
        "Taluka Head Direct Commission: 10% (₹14,000).",
        "Commission paid annually upon subscription renewal.",
        "Onboard 5 schools to generate ₹70,000 annually."
    ], "Present a realistic math example for a mid-sized school cloud ERP license deal.")

    add_content_slide(prs, "Revenue Sharing", "Worked Example: Smart Classroom Install", [
        "Client: Wagholi English School (6 classrooms).",
        "Smart Board Package: ₹1,20,000 per room.",
        "Total Hardware Invoice Value: ₹7,20,000.",
        "Taluka Head Direct Commission: 4% (₹28,800).",
        "Commission paid post-installation and sign-off check.",
        "Cross-sell opportunity: Annual Maintenance Contract (AMC)."
    ], "Present a realistic math example for a smart board hardware sourcing deal, showing cross-sell hooks.")

    add_content_slide(prs, "Revenue Sharing", "Monthly Performance Incentives", [
        "Milestone 1: ₹2L sales value in a month.",
        "  - Unlocks extra 1% bonus commission.",
        "Milestone 2: 10 active schools + 5 vendors/quarter.",
        "  - Unlocks flat ₹15,000 DTC travel override.",
        "Ensures administrative costs are covered during growth.",
        "Standardized structures apply across Pune district."
    ], "Detail the performance incentives that partners can unlock by hitting monthly and quarterly milestones.")

    # SECTION 8
    add_section_divider(prs, 8, "Product Catalogue")
    
    add_content_slide(prs, "Product Catalogue", "Software ERP & LMS Solutions", [
        "Problem: Paper registers waste school administrative time.",
        "Solution: Integrated cloud school management database ERP.",
        "Features: Fees tracking, automated mark sheets, WhatsApp alerts.",
        "Target Customer: Schools, senior colleges, coaching classes.",
        "Pricing Model: Annual subscription per student.",
        "Cross-sell: RFID biometric tracking gates."
    ], "Detail the school ERP product, explaining the problem, features, target market, and pricing model.")

    add_content_slide(prs, "Product Catalogue", "Hardware & Smart Board Packages", [
        "Problem: Traditional classrooms limit interactive learning.",
        "Solution: Full HD interactive touch smart screens.",
        "Features: Short-throw projectors, multi-touch screens, sound columns.",
        "Target Customer: Private schools, coaching academies, ITIs.",
        "Pricing Model: One-time hardware purchase fee.",
        "Cross-sell: Classroom audio columns, AMC packages."
    ], "Detail the Smart Classroom hardware package, highlighting components, pricing models, and cross-sell options.")

    add_content_slide(prs, "Product Catalogue", "B2B Marketplace Sourcing Portal", [
        "Problem: High markups from school uniforms distributors.",
        "Solution: Direct marketplace connecting schools with factories.",
        "Features: Pre-vetted catalog, transparent margin billing, logistics.",
        "Target Customer: School procurement heads, trust directors.",
        "Pricing Model: Bulk order transaction volume fee.",
        "Cross-sell: PE sports gear, custom student diaries."
    ], "Detail the DnyanMitra marketplace portal, showing how it aggregates bulk orders to bypass middlemen.")

    add_content_slide(prs, "Product Catalogue", "Institutional Excellence Audits", [
        "Problem: Trustees lack data on school tech gaps.",
        "Solution: Consultative AI-driven tech readiness assessments.",
        "Features: Safety check, administrative audit, pedagogical scoring.",
        "Target Customer: School management boards and trustees.",
        "Pricing Model: First audit free, follow-ups ₹5,000.",
        "Cross-sell: Smart boards and ERP subscriptions."
    ], "Outline the IX framework audit service, which acts as a key sales prospecting tool.")

    add_content_slide(prs, "Product Catalogue", "Faculty Training Services", [
        "Problem: School staff struggle to adopt ERPs.",
        "Solution: Hands-on user workshops led by DTC.",
        "Features: 2-day on-site training sessions, user manuals.",
        "Target Customer: School teachers and admin staff.",
        "Pricing Model: Flat training package fee (₹8,000).",
        "Cross-sell: Advanced AI exam grading modules."
    ], "Detail the training services offered by the DTC to ensure high software adoption rates.")

    add_content_slide(prs, "Product Catalogue", "Support & AMC Packages", [
        "Problem: Smart boards sit broken for weeks.",
        "Solution: Annual Maintenance Contracts with 24-hour SLA.",
        "Features: Quarterly preventive checks, priority hardware replacements.",
        "Target Customer: Schools with existing hardware setups.",
        "Pricing Model: Annual contract fee (12% commission).",
        "Cross-sell: Classroom projector bulb replacements."
    ], "Highlight the recurring AMC service line that provides steady income while keeping school hardware functional.")

    # SECTION 9
    add_section_divider(prs, 9, "Service Portfolio")
    
    add_content_slide(prs, "Service Portfolio", "Installation & Deployment Services", [
        "Professional mounting of smart board displays.",
        "Projector calibration and cabling checks.",
        "Local network mapping for classroom connections.",
        "Biometric attendance gateway testing and integration.",
        "Conducted by certified regional DTC engineers.",
        "SLA: Installation within 5 business days."
    ], "Detail the technical setup services that the DTC coordinates for schools post-purchase.")

    add_content_slide(prs, "Service Portfolio", "Training & Adoption Services", [
        "ERP system onboarding workshops for administrators.",
        "LMS dashboard training sessions for teachers.",
        "User guides, training videos, and documentation access.",
        "Follow-up checks 30 days post-launch.",
        "Maintains high user active rates locally.",
        "Reduces software support ticket volume."
    ], "Explain how the DTC handles user onboarding to ensure that teachers and staff adopt the cloud ERP systems.")

    add_content_slide(prs, "Service Portfolio", "Annual Maintenance Contracts (AMC)", [
        "Quarterly on-site preventive device checks.",
        "Free diagnostic testing for projector lamps.",
        "Standardized hardware repair ticketing workflow.",
        "Priority component shipping from Pune warehouse.",
        "Guarantees school device uptime throughout year.",
        "Creates high-margin recurring contract revenue."
    ], "Detail the recurring AMC portfolio, highlighting response SLAs, spare parts, and preventive dispatches.")

    add_content_slide(prs, "Service Portfolio", "Audit & Sourcing Advisory", [
        "Consultative Institutional Excellence auditing surveys.",
        "AI technology gap mapping calculations.",
        "Safety compliance mapping (CCTV coverage, biometrics).",
        "Local vendor audit site visits.",
        "Procurement budgeting workshops for school trustees.",
        "Establishes DTC as trusted regional advisor."
    ], "Summarize the audit and consulting services that position the DTC as a tech advisor rather than a sales agency.")

    # SECTION 10
    add_section_divider(prs, 10, "Sales Process")
    
    add_content_slide(prs, "Sales Process", "Lead Generation & Outreach", [
        "FSE maps schools and identifies trustees.",
        "Initial calling confirms right principal contact.",
        "Outreach uses Appointing wording rather than Recruiting.",
        "Curiosity-based openings avoid pricing talk early.",
        "Goal: Schedule first demo meeting slot.",
        "All leads logged inside DASP CRM."
    ], "Outline the initial stages of the sales process, focusing on mapping local schools and scheduling first meetings.")

    add_content_slide(prs, "Sales Process", "Demos & On-site Surveys", [
        "First Meeting: Build relevance and listen.",
        "Identify school pain points (CCTV, manual fees).",
        "Run free Institutional Excellence on-site survey.",
        "Map classroom dimensions for smart boards.",
        "FSE completes survey logs in CRM.",
        "Goal: Deliver customized digital readiness report."
    ], "Detail the audit and demo stage, showing how the free on-site survey acts as a consultative sales tool.")

    add_content_slide(prs, "Sales Process", "Proposals, Closing & Onboarding", [
        "Taluka Head drafts proposal and quote.",
        "Presents ROI cost-savings sheet to trustees.",
        "District Head joins high-value board meetings.",
        "DASP HQ handles invoicing post-signing.",
        "Technical team completes setup and training.",
        "DTC coordinates post-sale support and updates."
    ], "Detail the contract finalization, explaining proposal drafting, trustee negotiations, billing, and setup handover.")

    add_content_slide(prs, "Sales Process", "Customer Success & Renewals", [
        "FSE visits school monthly to check usage.",
        "Quarterly maintenance visits prevent device downtime.",
        "Helpdesk resolves billing/technical queries within SLA.",
        "Renewals mapped in CRM 60 days before expiration.",
        "AMC signups pitching starts in Month 10.",
        "Cross-sell campaigns trigger for next term."
    ], "Explain post-sale support, highlighting monthly visits, AMC pitching, and annual license renewals.")

    # SECTION 11
    add_section_divider(prs, 11, "12-Month Business Plan")
    
    # Month 1-3 roadmap
    m1_3 = [
        {"title": "Month 1: Setup", "desc": "Admin training, CRM setup, map 100% schools & 15 vendors.", "kpi": "DTC Setup Signoff"},
        {"title": "Month 2: Audits", "desc": "Recruit FSE, start visits, complete 20 school surveys.", "kpi": "FSE Hired & Active"},
        {"title": "Month 3: Launch", "desc": "Deploy ERP to early adopters, onboard 5 local vendors.", "kpi": "First Sales Logged"}
    ]
    add_roadmap_slide(prs, "12-Month Plan", "Months 1 to 3: Setup & Launch Phase", m1_3, "Walk the candidate through the first quarter, showing that the focus is on training, recruitment, school mapping, and initial sales.")

    # Month 4-6 roadmap
    m4_6 = [
        {"title": "Month 4: ERP", "desc": "Run ERP demos, host AI awareness, onboard 5 vendors.", "kpi": "5 ERP Deals Closed"},
        {"title": "Month 5: Drive", "desc": "Onboard 10 vendors, run procurement trustee meetings.", "kpi": "₹1.5L Sales Value"},
        {"title": "Month 6: Workshop", "desc": "Host first DTC workshop, complete 15 school audits.", "kpi": "15 IX Audits Signed"}
    ]
    add_roadmap_slide(prs, "12-Month Plan", "Months 4 to 6: Sourcing Expansion Phase", m4_6, "Explain the second quarter operations, focusing on bulk marketplace deals, vendor drives, and hosting the first regional workshop.")

    # Month 7-9 roadmap
    m7_9 = [
        {"title": "Month 7: Hardware", "desc": "Launch smart classroom and projector campaigns.", "kpi": "5 Smart Rooms Set Up"},
        {"title": "Month 8: Sports", "desc": "Run PE gear sourcing, onboard sports academies.", "kpi": "KridaMitra Launch"},
        {"title": "Month 9: AI", "desc": "Run automated grading and faculty AI workshops.", "kpi": "10 Schools Active"}
    ]
    add_roadmap_slide(prs, "12-Month Plan", "Months 7 to 9: Advanced Solutions Phase", m7_9, "Discuss the third quarter, showing how the DTC transitions into high-value smart classrooms, PE sports bookings, and AI modules.")

    # Month 10-12 roadmap
    m10_12 = [
        {"title": "Month 10: AMCs", "desc": "Initiate AMC pitching, run budget school meetings.", "kpi": "80% Accounts Covered"},
        {"title": "Month 11: Proposals", "desc": "Draft large multi-school proposals, plan expansion.", "kpi": "₹3L Pipeline Built"},
        {"title": "Month 12: Review", "desc": "Annual review, CSAT surveys, set Year 2 targets.", "kpi": "Year 2 Charter Signed"}
    ]
    add_roadmap_slide(prs, "12-Month Plan", "Months 10 to 12: Recurring Revenue Phase", m10_12, "Present the final quarter, focusing on AMC contract lock-ins, next-year budget planning, annual reviews, and Year 2 goals.")

    # SECTION 12
    add_section_divider(prs, 12, "Frequently Asked Questions")
    
    add_content_slide(prs, "FAQs", "Business Model & Territory FAQs", [
        "Q: What is the core business model? (A: DTC B2B franchise).",
        "Q: Is there an upfront fee? (A: No fee; partner invests in local DTC setup).",
        "Q: Are territory boundaries exclusive? (A: Yes, based on Maharashtra revenue block limits).",
        "Q: Can we sell outside Haveli? (A: No, territory limits are strictly enforced).",
        "Q: Can I manage two talukas? (A: Only if targets are consistently met in Haveli)."
    ], "Answer common questions about territory boundaries, exclusivity, and DTC setup investments.")

    add_content_slide(prs, "FAQs", "Financial & Payout FAQs", [
        "Q: What is the commission payout cycle? (A: Settled by the 10th of every month).",
        "Q: Is tax deducted at source? (A: Yes, TDS is applied as per government rules).",
        "Q: How is recurring revenue earned? (A: Annual ERP renewals and AMC contract sign-ups).",
        "Q: Are commissions paid on refunds? (A: Commissions are adjusted in the next cycle).",
        "Q: Can we offer custom discounts? (A: Only with written District Head approval)."
    ], "Answer queries related to monthly payment cycles, TDS deductions, recurring SaaS renewals, and discount policies.")

    add_content_slide(prs, "FAQs", "Operations & Support FAQs", [
        "Q: Who hires the local FSE? (A: The Taluka Head recruits and funds the FSE).",
        "Q: What is the installation SLA? (A: Smart boards mounted within 5 business days).",
        "Q: Who provides customer tech support? (A: DTC handles L1; DASP escalates to L2/HQ).",
        "Q: Are marketing print files free? (A: Yes, digital vector files are free).",
        "Q: What training do we receive? (A: 1-week mandatory online CRM/product training)."
    ], "Address queries about FSE recruitment, installation SLAs, technical ticket escalation, and marketing files.")

    add_content_slide(prs, "FAQs", "Legal & Future Growth FAQs", [
        "Q: What is the contract term? (A: 12 months, renewable annually based on KPIs).",
        "Q: What are grounds for termination? (A: Code violations, data leaks, or unverified vendors).",
        "Q: Do we represent future platforms? (A: Active partners get first right of refusal).",
        "Q: Who owns client data records? (A: All records remain DASP Digital property).",
        "Q: Can I cancel the agreement? (A: Yes, with a 30-day written notice)."
    ], "Provide answers regarding contract durations, non-compete rules, data ownership, and future platform launches.")

    # SECTION 13
    add_section_divider(prs, 13, "Risk Assessment")
    
    # Risk table
    risk_headers = ["Identified Sourcing Risk", "Potential Impact", "DASP Digital Mitigation Plan", "DTC Action Plan"]
    risk_rows = [
        ["Erratic rural electricity", "High (classroom downtime)", "Hardware bundled with hybrid UPS setups", "Verify school backup generators"],
        ["Delays in trustee board sign-off", "Medium (longer sales cycle)", "AI-generated IX audit report hooks", "Schedule regular face-to-face updates"],
        ["FSE turnover and attrition", "Medium (visit delays)", "Standardized FSE onboarding training", "Maintain pipeline of local CVs"],
        ["Cheap offline software competitors", "Low (price wars)", "Focus on local DTC support presence", "Conduct free computer lab tech audits"]
    ]
    add_table_slide(prs, "Risk Assessment", "DTC Risk Mitigation Matrix", risk_headers, risk_rows, "Present the local risks in semi-urban India and show how DASP's model mitigates them.")

    # SECTION 14
    add_section_divider(prs, 14, "Success Stories")
    
    add_content_slide(prs, "Success Stories", "Case Study: Smart Classroom Upgrade", [
        "Client: Private Secondary School (Haveli pilot).",
        "Installed: 5 Interactive Smart Boards + CRM link.",
        "Total Order Value: ₹6,00,000 (excl. tax).",
        "DTC Commission: ₹24,000 direct payout.",
        "Results: Class participation improved by 40% (CSAT audit).",
        "School subsequently signed up for annual software ERP."
    ], "Present a realistic, illustrative case study showing how a hardware smart board sale can lead to recurring ERP deals.")

    add_content_slide(prs, "Success Stories", "Case Study: ERP deployment", [
        "Client: Regional Junior College (750 students).",
        "Deployed: DnyanMitra cloud ERP & biometric gates.",
        "ERP Contract: ₹1,50,000 annual SaaS license.",
        "DTC Commission: ₹15,000 recurring payout.",
        "Results: Fee collection administration time reduced by 50%.",
        "First-line support ticket volume: 0 critical bugs."
    ], "Present a second case study detailing software ERP deployment and the resulting admin efficiency for the school.")

    # SECTION 15
    add_section_divider(prs, 15, "Partnership Journey")
    
    add_content_slide(prs, "Partnership Journey", "Vetting & Evaluation Stage", [
        "Meeting 1: First demo session (interest check).",
        "Meeting 2: Taluka Business Workshop (opportunity check).",
        "Vetting scorecard completed by District Head.",
        "Physical site visit of training centre.",
        "Bank details verified via Pennydrop test.",
        "Meeting 3: Final onboarding and agreement sign-off."
    ], "Explain the three-meeting selection process that leads up to the signing of the partnership contract.")

    add_content_slide(prs, "Partnership Journey", "Enablement & Growth Stage", [
        "Step 1: Complete 1-week coordinator training.",
        "Step 2: Provision CRM and email dashboard credentials.",
        "Step 3: Recruit Field Sales Executive.",
        "Step 4: Execute first month target school audits.",
        "Step 5: Hold weekly review meetings with District Head.",
        "Step 6: Plan Year 2 goals and DTC scaling."
    ], "Detail the steps after contract signing, focusing on software provisioning, FSE hiring, and scaling operations.")

    # SECTION 16
    add_section_divider(prs, 16, "Next Steps")
    
    add_content_slide(prs, "Next Steps", "Administrative Vetting Requirements", [
        "Provide PAN Card and Aadhaar Card uploads.",
        "Provide bank account details (mandate form).",
        "Provide GSTIN registration details (if active).",
        "Share physical address proof of training centre.",
        "Verify FSE recruit contact information.",
        "Submit signed pre-agreement checklist sheet."
    ], "Detail the documents the candidate needs to submit within the next 48 hours to finalize onboarding.")

    add_content_slide(prs, "Next Steps", "Launch Checklist (First 48 Hours)", [
        "1. Log into your custom CRM dashboard.",
        "2. Join official district coordinator WhatsApp groups.",
        "3. Review the DnyanMitra product catalogue.",
        "4. Print DTC banners and standees.",
        "5. Compile list of top 20 target schools.",
        "6. Schedule the first weekly call with District Head."
    ], "Provide the immediate launch checklist to guide the new partner's first 48 hours of business operations.")

    prs.save(file_path)
    print(f"Presentation generated successfully: {file_path}")

if __name__ == '__main__':
    dest_path = r"D:\company\products\dnyanmitra-knowledge-center\content\presentations\DM-PPT-Taluka-Head-Due-Diligence-Haveli-v1.0.pptx"
    # Ensure directory exists
    os.makedirs(os.path.dirname(dest_path), exist_ok=True)
    build_presentation(dest_path)
